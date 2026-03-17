import os
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_pptx():
    prs = Presentation()

    # Paths
    brain_dir = r"C:\Users\Administrator\.gemini\antigravity\brain\b34e3501-2ed3-4089-bfcd-c2274a685cc8"
    img_goals = os.path.join(brain_dir, "goal_hierarchy_diagram_1773341332729.png")
    img_caps = os.path.join(brain_dir, "capability_grouping_diagram_1773341348464.png")
    img_acq = os.path.join(brain_dir, "acquaintance_diagram_1773341363082.png")
    img_int = os.path.join(brain_dir, "interaction_diagram_1773341482787.png")

    def add_title_slide(title_text, subtitle_text, notes=""):
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = title_text
        subtitle.text = subtitle_text
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    def add_bullet_slide(title_text, bullets, notes=""):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = title_text
        tf = slide.placeholders[1].text_frame
        for bullet in bullets:
            p = tf.add_paragraph()
            if isinstance(bullet, tuple):
                p.text = bullet[0]
                p.level = bullet[1]
            else:
                p.text = bullet
                p.level = 0
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    def add_image_slide(title_text, image_path, notes=""):
        slide_layout = prs.slide_layouts[5] # Blank slide with title
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = title_text
        if os.path.exists(image_path):
            slide.shapes.add_picture(image_path, Inches(1), Inches(1.5), width=Inches(8))
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    # 1. Title Slide
    add_title_slide(
        "Crop Disease Intelligent Agent",
        "DCIT 403: Intelligent Agent Systems\nSemester Project Presentation\nUniversity of Ghana",
        "Welcome to the presentation for the Crop Disease Intelligent Agent. This project explores the intersection of Artificial Intelligence and Agriculture. We followed the Prometheus Methodology."
    )

    # 2. Problem Description
    add_bullet_slide(
        "1.1 Problem Description",
        [
            "The Challenge: Global food security threats.",
            "10-16% annual loss in global agricultural production.",
            "Billions of dollars lost; smallholders hit hardest.",
            "Critical Gap: Access to Plant Pathologists.",
            "The Solution: Edge AI for field diagnosis."
        ],
        "We address the lack of expert diagnostic capacity in rural areas. By the time a farmer notices a disease, the yield loss is irreversible. Our system provides a triage layer—allowing for rapid identification."
    )

    # 3. Why an Agent?
    add_bullet_slide(
        "1.2 Why an Intelligent Agent?",
        [
            "Agent Properties: Autonomy, Reactivity, Proactiveness.",
            "BDI Architecture: Beliefs, Desires, Intentions.",
            "Situated in an environment (the field).",
            "Goal-directed behavior."
        ],
        "An Intelligent Agent is situated in an environment. Our agent maintains 'Beliefs' about its own uncertainty and the history of the crop. It is proactive because it considers the severity and decides whether to escalate."
    )

    # 4. Phase 1: Goal Specification
    add_bullet_slide(
        "PHASE 1: Goal Specification",
        [
            "Mission: Protect Crop Health.",
            "G1: Diagnose Diseases.",
            "G2: Treatment Recommendations.",
            "G3: Field Monitoring."
        ],
        "We break down the mission into measurable sub-goals. This ensures the system is robust and that we have clear metrics for success."
    )

    # 5. Goal Hierarchy Diagram
    add_image_slide("Goal Hierarchy Diagram", img_goals, "Functional decomposition showing the system mission and its sub-goals.")

    # 6. Phase 2: Architectural Design
    add_bullet_slide(
        "PHASE 2: Architectural Design",
        [
            "Single-Agent BDI Orchestrator.",
            "Capabilities: Perception, Decision, Action.",
            "Separation of concerns."
        ],
        "We modularized the functionalities into three main capabilities. Perception handles how it sees (CNN), and Decision handles how it thinks."
    )

    # 7. Capability Grouping Diagram
    add_image_slide("Capability Grouping Diagram", img_caps, "How functionalities are grouped into high-level capability modules.")

    # 8. Phase 3: Interaction Design
    add_bullet_slide(
        "PHASE 3: Interaction Design",
        [
            "Scenario-based design.",
            "Protocols for farmer and environment interaction.",
            "Handling uncertainty."
        ],
        "Interaction design defines the protocols between entities. We designed for both success cases and edge cases like low confidence."
    )

    # 9. Sequence Interaction Diagram
    add_image_slide("Interaction Diagram", img_int, "The flow of messages between the farmer, the agent, and the environment.")

    # 10. Phase 4: Detailed Design
    add_bullet_slide(
        "PHASE 4: Detailed Design - Plans",
        [
            "BDI Plans: Diagnose, Escalate, Monitor.",
            "Plan selection based on beliefs.",
            "Threshold-driven decision logic."
        ],
        "Plans are recipes for achieving goals. The agent selects its plan based on its current 'Beliefs'."
    )

    # 11. Phase 5: Implementation
    add_bullet_slide(
        "PHASE 5: Implementation",
        [
            "Python, PyTorch (ResNet-18), Gradio.",
            "Custom BDI implementation.",
            "Environment simulation."
        ],
        "Technically, we used PyTorch for the CNN layer. ResNet-18 was chosen for accuracy and performance. The BDI logic was custom-built."
    )

    # 12. Conclusion
    add_bullet_slide("Conclusion", ["Autonomous identification achieved.", "Prometheus methodology validated.", "Future: IoT sensor integration."], "The system acts as an intelligent partner for farmers. It is a true Intelligent Agent.")

    # Save to Desktop with a timestamp to avoid locks
    desktop_path = os.path.join(os.environ['USERPROFILE'], 'Desktop')
    timestamp = int(time.time())
    output_path = os.path.join(desktop_path, f"Crop_Agent_Presentation_Detailed_{timestamp}.pptx")
    
    # Also try to save to the standard path but catch error
    std_path = os.path.join(desktop_path, "Crop_Disease_Agent_Project_Presentation.pptx")
    
    try:
        prs.save(std_path)
        print(f"Presentation saved to standard path: {std_path}")
    except Exception as e:
        print(f"Could not save to standard path (likely open). Error: {e}")
        prs.save(output_path)
        print(f"Presentation saved to alternative path: {output_path}")

if __name__ == "__main__":
    create_pptx()
